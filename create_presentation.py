"""
Create PowerPoint presentation from experiment results.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def create_title_slide(prs):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Multi-Modal Deep Learning for\nOrthodontic Malocclusion Classification"
    subtitle.text = "Comparing Three Experimental Approaches\n\nPhase 1 Experiments"

    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True

    return slide


def create_content_slide(prs, title_text, content_items):
    """Create a slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    # Add content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    return slide


def create_table_slide(prs, title_text, headers, rows):
    """Create a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5) * rows_count

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # Set headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Set data rows
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(14)

    return slide


def create_image_slide(prs, title_text, image_path, caption=None):
    """Create a slide with an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    # Add image
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)

    slide.shapes.add_picture(str(image_path), left, top, width=width)

    # Add caption if provided
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].font.size = Pt(14)
        caption_frame.paragraphs[0].font.italic = True

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating PowerPoint presentation...")

    # Slide 1: Title
    print("  [1/10] Title slide")
    create_title_slide(prs)

    # Slide 2: Problem Statement
    print("  [2/10] Problem statement")
    create_content_slide(prs, "Problem Statement", [
        "Objective: Three-class malocclusion classification (Class I, II, III)",
        "",
        "Challenge: How to effectively leverage multi-modal medical imaging?",
        "",
        "Image Modalities:",
        "  • OPG (Panoramic X-rays)",
        "  • Intraoral photographs",
        "  • Cephalometric X-rays",
        "",
        "Dataset: 422 unique patients with patient-aware splitting (70/15/15)"
    ])

    # Slide 3: Experimental Approaches
    print("  [3/10] Experimental approaches")
    create_table_slide(prs, "Experimental Approaches",
        ["Experiment", "Approach", "Architecture"],
        [
            ["E1", "OPG-only Baseline", "MedGemma-4B (frozen) + classifier"],
            ["E2", "Naïve Late Fusion", "3 CNNs (one per modality) + ensemble"],
            ["E3", "Multi-Image Prompting", "MedGemma-4B with multi-image tokens"]
        ]
    )

    # Slide 4: Experiment 1 Details
    print("  [4/10] Experiment 1 details")
    create_content_slide(prs, "Experiment 1: OPG-only Baseline", [
        "Strategy: Establish single-modality baseline using panoramic X-rays",
        "",
        "Architecture:",
        "  • Frozen MedGemma-4B vision encoder (1152-dim features)",
        "  • Trainable classification head",
        "  • Robust 3-tier modality filtering",
        "",
        "Training:",
        "  • 20 epochs, batch size 8",
        "  • AdamW optimizer, Cosine LR scheduler",
        "  • Only OPG images used",
        "",
        "Best Performance: 66.67% validation accuracy (epoch 1)"
    ])

    # Slide 5: Experiment 2 Details
    print("  [5/10] Experiment 2 details")
    create_content_slide(prs, "Experiment 2: Naïve Late Fusion", [
        "Strategy: Train separate classifiers per modality, fuse at decision level",
        "",
        "Architecture:",
        "  • Intraoral CNN: SimpleClassifier (512-dim) → intraoral images",
        "  • OPG CNN: SimpleClassifier (512-dim) → panoramic X-rays",
        "  • Ceph CNN: SimpleClassifier (512-dim) → cephalometric X-rays",
        "  • Fusion: Average predictions from all three classifiers",
        "",
        "Training: Sequential training, 45 epochs each modality",
        "",
        "Best Performance: 93.51% validation accuracy ⭐"
    ])

    # Slide 6: Experiment 3 Details
    print("  [6/10] Experiment 3 details")
    create_content_slide(prs, "Experiment 3: Multi-Image Prompting", [
        "Strategy: Feed multiple images with specialized prompts to MedGemma",
        "",
        "Architecture:",
        "  • Frozen MedGemma-4B encoder + classification head",
        "  • Multi-image input with <start_of_image> tokens",
        "  • Modality-aware prompt engineering",
        "",
        "Prompt Example:",
        '  "<start_of_image> <start_of_image> <start_of_image>',
        "   You are analyzing multiple orthodontic images...",
        "   Image 1 is a panoramic X-ray (OPG).",
        '   Image 2 is an intraoral photograph..."',
        "",
        "Best Performance: 72.73% validation accuracy"
    ])

    # Slide 7: Results Comparison
    print("  [7/10] Results comparison")
    create_image_slide(prs, "Results Comparison",
        Path("visualizations/all_experiments_comparison.png"),
        "Comparison of training and validation metrics across all experiments")

    # Slide 8: E1 Training Dynamics
    print("  [8/10] E1 training dynamics")
    create_image_slide(prs, "E1: OPG-only Baseline - Training Dynamics",
        Path("visualizations/E1_OPG-only_Baseline_metrics.png"),
        "Shows overfitting - validation loss increases while training continues")

    # Slide 9: E2 Training Dynamics
    print("  [9/10] E2 training dynamics")
    create_image_slide(prs, "E2: Late Fusion - Training Dynamics (Best Performer)",
        Path("visualizations/E2_Late_Fusion_metrics.png"),
        "Extremely stable - maintains ~93.5% validation accuracy throughout")

    # Slide 10: Key Findings
    print("  [10/10] Key findings")
    create_content_slide(prs, "Key Findings & Conclusions", [
        "Performance Ranking:",
        "  1. E2 (Late Fusion): 93.51% ⭐ - Clear winner",
        "  2. E3 (Multi-Image): 72.73% - Moderate improvement over baseline",
        "  3. E1 (OPG-only): 66.67% - Baseline",
        "",
        "Insights:",
        "  ✓ Multi-modal fusion significantly outperforms single modality",
        "  ✓ Simple late fusion beats sophisticated multi-image prompting",
        "  ⚠ E3 shows promise but may need more training or tuning",
        "  ⚠ E1 exhibits overfitting with frozen MedGemma encoder",
        "",
        "Recommendation: Deploy E2 (Late Fusion) for production use"
    ])

    # Save presentation
    output_path = Path("presentation.pptx")
    prs.save(str(output_path))
    print(f"\n✓ Presentation saved to: {output_path}")
    print(f"  File size: {output_path.stat().st_size / 1024 / 1024:.2f} MB")


if __name__ == '__main__':
    main()
